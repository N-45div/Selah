from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from ..models import ServicePlan


def generate_pptx_deck(plan: ServicePlan) -> BytesIO:
    """
    Generates a 16:9 widescreen PowerPoint presentation (.pptx)
    with deep dark church booth background, high-contrast serif typography,
    and phonetic transliteration lines.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)  # 16:9 Widescreen standard
    blank_layout = prs.slide_layouts[6]

    # Title Slide
    title_slide = prs.slides.add_slide(blank_layout)
    bg = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(14, 13, 11)
    bg.line.fill.background()

    tx = title_slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(3))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = plan.service_name
    p.font.name = "Georgia"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = f"{plan.stream_title} • Selah Telecast"
    p2.font.name = "Calibri"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(184, 178, 165)
    p2.alignment = PP_ALIGN.CENTER

    # Song Slides
    for song in plan.songs:
        for s in song.slides:
            slide = prs.slides.add_slide(blank_layout)

            # Dark Background
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(14, 13, 11)
            bg.line.fill.background()

            # Header info (Song title & section)
            hdr = slide.shapes.add_textbox(Inches(1), Inches(0.6), Inches(11.333), Inches(0.8))
            htf = hdr.text_frame
            hp = htf.paragraphs[0]
            hp.text = f"{song.title} • {s.label}"
            hp.font.name = "Calibri"
            hp.font.size = Pt(14)
            hp.font.bold = True
            hp.font.color.rgb = RGBColor(212, 145, 42)  # Amber
            hp.alignment = PP_ALIGN.LEFT

            # Lyrics textbox
            tx = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11.333), Inches(5.0))
            tf = tx.text_frame
            tf.word_wrap = True

            for idx, line in enumerate(s.lines):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.text = line
                p.font.name = "Georgia"
                p.font.size = Pt(36)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER

                # Transliteration line
                if s.transliteration and idx < len(s.transliteration):
                    tp = tf.add_paragraph()
                    tp.text = s.transliteration[idx]
                    tp.font.name = "Calibri"
                    tp.font.italic = True
                    tp.font.size = Pt(22)
                    tp.font.color.rgb = RGBColor(240, 197, 110)  # Amber transliteration
                    tp.alignment = PP_ALIGN.CENTER

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
